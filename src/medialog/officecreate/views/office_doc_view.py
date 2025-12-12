# -*- coding: utf-8 -*-

# from medialog.officecreate import _
from Products.Five.browser import BrowserView
from zope.interface import implementer
from zope.interface import Interface
from plone.namedfile.file import NamedBlobImage
from plone.app.textfield.value import RichTextValue
from html2docx import html2docx
from plone.app.querystring import queryparser
#from Products.CMFCore.utils import getToolByName
from plone.dexterity.utils import iterSchemata
from zope.schema import getFieldsInOrder
from plone import api

# docx
from docx import Document
from docx.shared import Mm
from docx.text.paragraph import Paragraph
from docxtpl import DocxTemplate, InlineImage, RichText
from lxml import etree

#For PowerPoint:
from pptx import Presentation

# Images / resizing
from PIL import Image

import tempfile
from io import BytesIO
import subprocess
# import os


EMU_PER_MM = 36000

class IOfficeDocView(Interface):
    """ Marker Interface for IOfficeDocView"""


    

# -------------------------------------------------------------
# EXTRACT PLACEHOLDER SIZES FROM WORD FILE
# -------------------------------------------------------------
def get_box_sizes_from_docx(template_docx):
        """
        Returns: dict placeholder -> (width_mm, height_mm)
        Returns empty dict if doc or doc.docx is invalid.
        """
         # Namespaces
        w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        wps_ns = "http://schemas.microsoft.com/office/word/2010/wordprocessingShape"
        root = template_docx.element
        ns = root.nsmap or {}

        results = {}
        for t in template_docx.element.iter("{%s}t" % w_ns):
            if not t.text:
                continue
            text = t.text.strip()
            if not (text.startswith("{{") and text.endswith("}}")):
                continue
            placeholder = text.strip("{} ").strip()
            # climb parents to find <a:graphic>
            parent = t.getparent()
            ext_cx = ext_cy = None
            while parent is not None:
                for graphic in parent.iter("{%s}graphic" % a_ns):
                    # inside <a:graphicData>
                    for gdata in graphic.iter("{%s}graphicData" % a_ns):
                        for wsp in gdata.iter("{%s}wsp" % wps_ns):
                            for spPr in wsp.iter("{%s}spPr" % wps_ns):
                                for xfrm in spPr.iter("{%s}xfrm" % a_ns):
                                    for ext in xfrm.iter("{%s}ext" % a_ns):
                                        try:
                                            ext_cx = int(ext.get("cx"))
                                            ext_cy = int(ext.get("cy"))
                                        except (TypeError, ValueError):
                                            continue
                if ext_cx is not None and ext_cy is not None:
                    break
                parent = parent.getparent()

            if ext_cx is not None and ext_cy is not None:
                width_mm = ext_cx / EMU_PER_MM
                height_mm = ext_cy / EMU_PER_MM
                results[placeholder] = (width_mm, height_mm)

        return results

    
    

@implementer(IOfficeDocView)
class OfficeDocView(BrowserView):
    # If you want to define a template here, please remove the template from
    # the configure.zcml registration of this view.
    # template = ViewPageTemplateFile('office_doc_view.pt')

    def __call__(self):
        context = self.context
        selected_file_id = self.request.form.get('selected_file', None)
        format_ = self.request.get("format", "docx")  # ?format=pdf or ?format=docx

        if not selected_file_id:
            return self.index()

        filen = api.content.get(UID=selected_file_id)
        if not filen or not hasattr(filen, 'file'):
            return "Invalid file selected."

        file_data = filen.file.data
        file_stream = BytesIO(file_data)
        file_stream.seek(0)

        # Determine file type
        portal_type = getattr(filen, 'portal_type', '').lower()

        # Handle DOCX (Word)
        filenavn = filen.file.filename
        if portal_type in ('file', 'document') and filenavn.endswith('.docx'):
            # Build a python-docx Document for inspecting boxes
            template_docx = Document(file_stream)
            box_sizes = get_box_sizes_from_docx(template_docx)
                    
            # --- 1. Get Text/Simple Replacements ---
            # Assuming get_doc_replacements() returns a dictionary of simple variables
            doc = DocxTemplate(file_stream)
            replacements = self.get_doc_replacements(context, doc)
            
            # --- 2. Add Image Replacements to the SAME Dictionary ---
            for schema in iterSchemata(context):
                schema_obj = schema(context)
                for name, field in getFieldsInOrder(schema):
                    value = getattr(schema_obj, name, None)
                    
                    if isinstance(value, NamedBlobImage) and getattr(value, 'data', None):
                        image_stream = BytesIO(value.data)
                        width = value._width # 1200
                        height = value._height # 800
                        sizes = box_sizes.get(name)
                        if sizes: 
                            image_width = sizes[0] # 400
                            image_height = sizes[1] # 440
                            ratio_w = width / image_width # 3
                            ratio_h = height / image_height # 2
                            ratio = min(ratio_w, ratio_h)
                            width = width / ratio
                            height = height / ratio
                            # TO DO: Find heighest percentace, check image with PIL and center image
                            replacements[name] = InlineImage(doc, image_stream, width=Mm(width), height=Mm(height))
                        else:
                            # Add the InlineImage object to the main replacements dictionary.
                            # The 'name' variable here is assumed to be the FIELD NAME in Plone, 
                            # which should match the JINJA2 TAG in your DOCX template (e.g., {{ name }} )
                            replacements[name] = InlineImage(doc, image_stream)

            # --- 3. RENDER ONCE with the complete dictionary ---
            doc.render(replacements)

            # --- 4. Save and return the file (Plone specific handling) ---
            # Use a BytesIO object to capture the o utput for Plone response
            output_stream = BytesIO()
            doc.save(output_stream)
            output_stream.seek(0)
            attachment_name = f"{filenavn}.docx"
            
            if format_ == "pdf":
                # --- 2) Write DOCX to temp file ---
                docx_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
                docx_temp.write(output_stream.getvalue())
                docx_temp.close()

                # --- 3) Convert DOCX → PDF using LibreOffice ---
                subprocess.run([
                    "libreoffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", "/tmp",
                    docx_temp.name
                ], check=True)

                pdf_path = docx_temp.name.replace(".docx", ".pdf")

                # --- 4) Return PDF ---
                with open(pdf_path, "rb") as f:
                    pdf_data = f.read()

                attachment_name = f"{filenavn}.pdf"
                self.request.response.setHeader("Content-Type", "application/pdf")

                self.request.response.setHeader(
                    "Content-Disposition",
                    f'attachment; filename="{attachment_name}"'
                )
                return pdf_data

            else:

                self.request.response.setHeader(
                    'Content-Type',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                )
                self.request.response.setHeader(
                    'Content-Disposition',
                    f'attachment; filename="{attachment_name}"'
                )

                return output_stream.getvalue() 
        


        # Handle PPTX (PowerPoint)
        elif portal_type in ('file', 'presentation') or filen.file.filename.endswith('.pptx'):
            prs = Presentation(file_stream)
            replacements = self.get_ppt_replacements(context, prs)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    
                    if shape.shape_type == 13:  # 13 means it's a picture
                        # Use image alt text for key
                        key = shape.image._filename.split('.')[0]
                        
                        if  hasattr(context, key):
                            img_obj = getattr(context, key, None)
                            
                            # must have .data                            
                            if getattr(img_obj, "data", None):
                                image_stream = BytesIO(img_obj.data)
                                
                                # Placeholder
                                left = shape.left
                                top = shape.top
                                width = shape.width
                                height = shape.height

                                # Load image
                                image_stream.seek(0)
                                pil_img = Image.open(image_stream)
                                img_w, img_h = pil_img.size  # px

                                # Convert EMU → px
                                EMU_PER_INCH = 914400
                                PX_PER_INCH = 96
                                frame_w_px = width / EMU_PER_INCH * PX_PER_INCH
                                frame_h_px = height / EMU_PER_INCH * PX_PER_INCH

                                # Scale to COVER
                                scale = max(frame_w_px / img_w, frame_h_px / img_h)
                                new_w_px = img_w * scale
                                new_h_px = img_h * scale

                                # Convert back to EMU
                                new_width = int(new_w_px / PX_PER_INCH * EMU_PER_INCH)
                                new_height = int(new_h_px / PX_PER_INCH * EMU_PER_INCH)

                                # Remove old picture
                                sp = shape._element
                                sp.getparent().remove(sp)

                                image_stream.seek(0)

                                # Add new picture at full scaled size (cover)
                                new_shape = shape.part.slide.shapes.add_picture(
                                    image_stream,
                                    left,
                                    top,
                                    width=width,
                                    height=height
                                )

                                # ---- Crop to match original placeholder ----
                                # Calculate crop ratios (0-1)
                                crop_left = max((new_width - width) / 2 / new_width, 0)
                                crop_top  = max((new_height - height) / 2 / new_height, 0)
                                crop_right = crop_left
                                crop_bottom = crop_top

                                # Apply cropping
                                new_shape.crop_left = crop_left
                                new_shape.crop_top = crop_top
                                new_shape.crop_right = crop_right
                                new_shape.crop_bottom = crop_bottom
                    
                    if not shape.has_text_frame:
                        continue

                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                                original_text = run.text
                                new_text = original_text
                                for key, value in replacements.items():
                                    # import pdb; pdb.set_trace()
                                    strvalue = str(value).replace("\r", "")
                                    #.replace("\n", "*")
                                    new_text = new_text.replace(key, strvalue)
                                    
                                if new_text != original_text:
                                    run.text = new_text
                            
            output_stream = BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)

            self.request.response.setHeader(
                'Content-Type',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            self.request.response.setHeader(
                'Content-Disposition',
                'attachment; filename="generated.pptx"'
            )

            return output_stream.getvalue()

        return "Unsupported or invalid file type."
    
    
    def get_doc_replacements(self, context, doc):
        """Get the replacements for the Word template."""
        replacements = {}        
        
        for schema in iterSchemata(context):
            schema_obj = schema(context)  # This adapts context to the schema
            for name, field in getFieldsInOrder(schema):
                value = getattr(schema_obj, name, None)
                if not value:
                    continue
                if isinstance(value, NamedBlobImage) and getattr(value, 'data', None):
                    continue                           
                            
                # if isinstance(value, RichTextValue) and getattr(value, 'output', None):
                #     html = value.raw
                #     subdoc_stream = BytesIO()
                #     #document = Document()
                #     document = Document()
                #     html2docx(html, document)
                #     document.save(subdoc_stream)
                #     subdoc_stream.seek(0)
                #     subdoc = Document().new_subdoc(subdoc_stream)
                #     replacements[name] = subdoc
                else:
                    # Normal field
                    replacements[name] = value
                    
        replacements['title'] = context.Title()
        replacements['description'] = context.Description()
        
        return replacements
    
    def get_ppt_replacements(self, context, doc):
        """Get the replacements for the PPT template."""
        replacements = {}        
        
        for schema in iterSchemata(context):
            schema_obj = schema(context)  # This adapts context to the schema
            for name, field in getFieldsInOrder(schema):
                value = getattr(schema_obj, name, None)
                replacements[name] = value
                    
        replacements['title'] = context.Title()
        replacements['description'] = context.Description()
        
        return replacements
        
    def find_docx_in_templates(self):
        # Define the path to search within (e.g., /templates)
        portal = api.portal.get()
        docx_items = []
        if portal.get('templates', False):            
            templates = portal.get('templates')
            
            # Search for all File content items within the folder
            # To do: search for 'template' / subject
            documents = api.content.find(
                portal_type='File',
                context=templates,
                
            )
            
            # Collect the items
            for document in documents: 
                obj = document.getObject()
                file_type = obj.file.contentType
                if file_type in  [
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                    ]:
                    docx_items.append({
                        'title': obj.Title(),
                        'file_type': file_type.split('.')[-1],
                        # 'url': obj.absolute_url, 
                        # 'object': obj,
                        'uuid': obj.UID(),
                        'file_name': obj.file.filename  # Include the filename as well for easier reference
                    })
        
        return docx_items
        
        
        
    