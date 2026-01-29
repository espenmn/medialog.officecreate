# -*- coding: utf-8 -*-

# from medialog.officecreate import _
from Products.Five.browser import BrowserView
from zope.interface import implementer
from zope.interface import Interface

from pptx import Presentation
from io import BytesIO
from ZPublisher.Iterators import filestream_iterator
# from Products.Five.browser.pagetemplatefile import ViewPageTemplateFile

class IPPMergeView(Interface):
    """ Marker Interface for IPPMergeView"""


@implementer(IPPMergeView)
class PPMergeView(BrowserView):
    # If you want to define a template here, please remove the template from
    # the configure.zcml registration of this view.
    # template = ViewPageTemplateFile('p_p_merge_view.pt')

    def __call__(self):
        request = self.request

        if "form.submit" not in request:
            return self.index()

        # PP1 → will be appended
        source_file = request.form.get("ppt1")

        # PP2 → base document
        target_file = request.form.get("ppt2")

        if not source_file or not target_file:
            return "Both files are required"

        source_ppt = Presentation(source_file)
        target_ppt = Presentation(target_file)

        self.append_presentation(source_ppt, target_ppt)

        output = BytesIO()
        target_ppt.save(output)
        output.seek(0)

        response = request.response
        response.setHeader(
            "Content-Type",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        response.setHeader(
            "Content-Disposition",
            'attachment; filename="pp2_with_pp1_appended.pptx"'
        )

        response.setBody(output.getvalue())
        return response
    
    
    def append_presentation(self, source_ppt, target_ppt):
        for slide in source_ppt.slides:
            new_slide = target_ppt.slides.add_slide(
                target_ppt.slide_layouts[6]  # blank
            )
            for shape in slide.shapes:
                new_slide.shapes._spTree.insert_element_before(
                    shape.element, 'p:extLst'
                )

